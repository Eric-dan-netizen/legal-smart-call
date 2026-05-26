#!/usr/bin/env python3
"""Convert slides.html to PPTX by screenshotting each slide via Playwright."""

import subprocess, os, json, tempfile, sys
from pathlib import Path

html_path = '/Users/apple/Documents/智能外呼系统/slides.html'
output_path = '/Users/apple/Documents/智能外呼系统/智能外呼系统_商业分析.pptx'
img_dir = Path(tempfile.mkdtemp())

print(f'HTML: {html_path}')
print(f'Output: {output_path}')
print(f'Temp images: {img_dir}')

# Step 1: Use Playwright to capture each slide as PNG
js_script = f"""
const {{ chromium }} = require('/tmp/pptx-gen/node_modules/playwright');
(async () => {{
  const browser = await chromium.launch({{ channel: 'chrome', headless: true }});
  const page = await browser.newPage({{
    viewport: {{ width: 1280, height: 720 }},
    deviceScaleFactor: 2,
  }});
  await page.goto('file://{html_path}', {{ waitUntil: 'networkidle' }});

  const slides = await page.locator('.slide').all();
  console.log(`Found ${{slides.length}} slides`);

  for (let i = 0; i < slides.length; i++) {{
    const filepath = '{img_dir}/slide_' + String(i+1).padStart(2,'0') + '.png';
    await slides[i].screenshot({{ path: filepath }});
    console.log(`Captured slide ${{i+1}} → ${{filepath}}`);
  }}

  await browser.close();
  console.log('Done capturing.');
}})();
"""

js_file = img_dir / 'capture.js'
js_file.write_text(js_script)

print('Launching Playwright...')
result = subprocess.run(
    ['node', str(js_file)],
    capture_output=True, text=True,
    cwd=img_dir,
    timeout=120,
)
print(result.stdout)
if result.returncode != 0:
    print('STDERR:', result.stderr)
    sys.exit(1)

# Step 2: Embed screenshots into PPTX
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

pngs = sorted(img_dir.glob('slide_*.png'))
print(f'\nEmbedding {len(pngs)} images into PPTX...')

for png in pngs:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(
        str(png),
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )

# Make first slide active
prs.save(output_path)
print(f'\nSaved: {output_path}')
print(f'Slides: {len(prs.slides)}')

# Cleanup
import shutil
shutil.rmtree(img_dir)
