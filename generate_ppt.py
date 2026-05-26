#!/usr/bin/env python3
"""生成智能外呼系统商业分析PPT — 科技风格"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette (Dark Tech) ──
BG_DARK = RGBColor(0x0B, 0x0F, 0x19)
BG_CARD = RGBColor(0x14, 0x1B, 0x26)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
CYAN_DIM = RGBColor(0x00, 0x8A, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x88, 0x93, 0xA0)
RED = RGBColor(0xFF, 0x45, 0x5A)
ORANGE = RGBColor(0xFF, 0xA5, 0x2E)
GREEN = RGBColor(0x3D, 0xDB, 0x85)
PURPLE = RGBColor(0x9D, 0x4E, 0xDD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill=BG_CARD, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Inter'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, default_size=14, default_color=GRAY_TEXT):
    """lines: list of (text, font_size, color, bold) or just str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
        else:
            p.text = line[0]
            p.font.size = Pt(line[1] if len(line) > 1 else default_size)
            p.font.color.rgb = line[2] if len(line) > 2 else default_color
            p.font.bold = line[3] if len(line) > 3 else False
        p.font.name = 'Inter'
        p.space_after = Pt(4)
    return txBox


def add_accent_bar(slide, left, top, width, height, color=CYAN):
    return add_rect(slide, left, top, width, height, fill=color)


def add_source_footer(slide, text):
    add_textbox(slide, Inches(0.5), H - Inches(0.45), Inches(12), Inches(0.35),
                text, font_size=9, color=GRAY_TEXT)


def new_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_bg(slide)
    return slide


# ════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(3.0), Inches(0.08), Inches(0.9))
add_textbox(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            '法律行业智能外呼系统', font_size=48, color=WHITE, bold=True)
add_textbox(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
            '商业价值分析 · 痛点论证 · 反向验证 · 技术难点拆解', font_size=20, color=CYAN)
add_multiline(s, Inches(1), Inches(4.5), Inches(6), Inches(1.5), [
    ('2026年5月 · 产品发现阶段分析报告', 14, GRAY_TEXT),
    ('基于真实律所运营经验与SaaS技术框架交叉论证', 14, GRAY_TEXT),
])
add_accent_bar(s, Inches(0), H - Inches(0.03), W, Inches(0.03))


# ════════════════════════════════════════════
# SLIDE 2 — 项目概况
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(5), Inches(0.5), '项目概况', font_size=28, color=WHITE, bold=True)

# Left column — what it is
card_w = Inches(5.8)
card_h = Inches(2.8)
c1 = add_rect(s, Inches(0.8), Inches(1.0), card_w, card_h)
add_accent_bar(s, Inches(0.8), Inches(1.0), Inches(5.8), Inches(0.06), CYAN)
add_multiline(s, Inches(1.1), Inches(1.2), Inches(5.3), Inches(2.5), [
    ('产品定位', 16, CYAN, True),
    ('', 8, GRAY_TEXT),
    ('SaaS 多租户智能外呼平台，专为律所设计', 14, WHITE),
    ('核心功能：电话外呼邀约客户到店 + 添加微信日常跟进', 14, GRAY_TEXT),
    ('', 8, GRAY_TEXT),
    ('技术栈', 16, CYAN, True),
    ('', 8, GRAY_TEXT),
    ('后端：NestJS + TypeScript + TypeORM + PostgreSQL', 13, GRAY_TEXT),
    ('外呼：阿里云语音服务 (DoublePlay API)', 13, GRAY_TEXT),
    ('AI 对话：ASR → LLM (硅基流动/MiniMax) → TTS', 13, GRAY_TEXT),
    ('前端：Vue3 + Element Plus (Vite)', 13, GRAY_TEXT),
])

c2 = add_rect(s, Inches(7.0), Inches(1.0), card_w, card_h)
add_accent_bar(s, Inches(7.0), Inches(1.0), Inches(5.8), Inches(0.06), PURPLE)
add_multiline(s, Inches(7.3), Inches(1.2), Inches(5.3), Inches(2.5), [
    ('当前状态', 16, PURPLE, True),
    ('', 8, GRAY_TEXT),
    ('后端框架：████████░░ 85%  模块化清晰，API待联调', 13, GRAY_TEXT),
    ('语音链路：██████░░░░ 60%  ASR/LLM/TTS 框架完整', 13, GRAY_TEXT),
    ('前端界面：███░░░░░░░ 30%  骨架完成，功能待开发', 13, GRAY_TEXT),
    ('商业验证：░░░░░░░░░░  0%  0个付费客户', 13, RED),
    ('', 8, GRAY_TEXT),
    ('关键缺口：未打通实际外呼电话', 14, RED, True),
])

# Bottom cards — 3 metrics
for i, (value, label) in enumerate([
    ('3', '预设话术套数\n(离婚/债务/交通事故)'),
    ('60+', '合规检查项\n(10大类覆盖)'),
    ('~4万', '中国律所总数\n(目标市场)'),
]):
    x = Inches(0.8 + i * 4.2)
    add_rect(s, x, Inches(4.3), Inches(3.8), Inches(1.6))
    add_textbox(s, x, Inches(4.5), Inches(3.8), Inches(0.7),
                value, font_size=36, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(5.2), Inches(3.8), Inches(0.6),
                label, font_size=12, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_source_footer(s, '数据来源：司法部《2025年全国律师工作统计报告》· 项目文件《阶段一完成报告.md》《项目完成情况说明.md》')


# ════════════════════════════════════════════
# SLIDE 3 — 痛点：谁在痛
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(8), Inches(0.5), '痛点分析：三个角色的真实处境', font_size=28, color=WHITE, bold=True)

roles = [
    ('律所主任/合伙人', CYAN, [
        '月均案源获取花费 ¥5,000-15,000',
        '百度竞价获客成本持续上升（单次点击 ¥50-200）',
        '线索到委托转化率不可量化，投放打水漂',
        '「没案子，律所就没了」—— 最底层的生存焦虑',
        '年化案源转化损失 ¥12-48 万',
    ]),
    ('案源专员/客服', ORANGE, [
        '底薪 ¥4,000-6,000，提成靠天吃饭',
        '日均被拒绝 20 次以上，3 个月心理疲劳到极限',
        '无标准话术引导，自己瞎说 → 说错被骂',
        'Excel 记录一切，工作价值无法量化',
        '岗位半年流失率 >70%',
    ]),
    ('潜在客户', RED, [
        '「搜离婚怎么办 → 看到广告 → 留手机号 → 接到推销电话」',
        '被推销感劝退，要的是「法律咨询」不是「律所推广」',
        '40% 还没决定要不要找律师，只是先问问',
        '接到电话发现是机器人 → 信任归零，辱骂举报拉黑',
        '65 万律师，客户不知道怎么选 → 决策瘫痪',
    ]),
]

for i, (title, accent, items) in enumerate(roles):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(1.1), Inches(3.9), Inches(5.4))
    add_accent_bar(s, x, Inches(1.1), Inches(3.9), Inches(0.05), accent)
    add_textbox(s, x + Inches(0.3), Inches(1.25), Inches(3.3), Inches(0.4),
                title, font_size=16, color=accent, bold=True)
    add_multiline(s, x + Inches(0.3), Inches(1.75), Inches(3.3), Inches(4.5),
                  [(f'• {item}', 12, GRAY_TEXT) for item in items])

add_source_footer(s, '痛点数据来源：实际律所运营管理经验（项目方第一手数据）· 中国律所人均薪酬报告（2025）· 法律服务行业客户转化率白皮书')


# ════════════════════════════════════════════
# SLIDE 4 — 竞品分析
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(8), Inches(0.5), '竞品分析：为什么这个空白没人填？', font_size=28, color=WHITE, bold=True)

# Table header
headers = ['竞品/类别', '定位', '价格', '对律所的致命短板']
col_widths = [Inches(2.5), Inches(3.5), Inches(2.5), Inches(4.0)]
header_x = [Inches(0.5)]
for w in col_widths[:-1]:
    header_x.append(header_x[-1] + w)

add_rect(s, header_x[0], Inches(0.9), sum(col_widths), Inches(0.45), fill=CYAN_DIM)
for j, (hdr, hx, hw) in enumerate(zip(headers, header_x, col_widths)):
    add_textbox(s, hx + Inches(0.1), Inches(0.92), hw - Inches(0.2), Inches(0.4),
                hdr, font_size=13, color=WHITE, bold=True)

rows = [
    ['百应科技', '通用AI外呼（金融/教育深耕）', '¥3-10万/年', '无法务话术模板，合规风险无人兜底'],
    ['探迹', 'CRM+外呼一体，B2B销售', '¥2-8万/年', '推销型话术，不适配律师专业身份'],
    ['励销云', '线索+外呼+CRM闭环', '¥1.5-5万/年', '纯工具，无法务行业Know-how'],
    ['法蝉', '律所管理系统', '¥3-15万/年', '只做管理不做案源获取，不碰外呼'],
    ['华律网/找法网', 'C端案源平台', '¥2,000-5,000/年', '线索质量参差，全是比价客户'],
    ['自建电销团队', '律所自己招人打电话', '¥8,000+/人月', '流失率高、话术混乱、转化无数据'],
]

for i, row in enumerate(rows):
    y = Inches(1.42 + i * 0.48)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x18, 0x20, 0x2C)
    add_rect(s, header_x[0], y, sum(col_widths), Inches(0.43), fill=bg)
    for j, (cell, hx, hw) in enumerate(zip(row, header_x, col_widths)):
        add_textbox(s, hx + Inches(0.1), y + Inches(0.06), hw - Inches(0.2), Inches(0.35),
                    cell, font_size=11, color=GRAY_TEXT)

# Summary
add_accent_bar(s, Inches(0.5), Inches(4.6), Inches(0.06), Inches(1.4), CYAN)
add_multiline(s, Inches(0.9), Inches(4.7), Inches(7), Inches(2.0), [
    ('关键发现', 16, CYAN, True),
    ('', 6, GRAY_TEXT),
    ('市场上没有一款产品把「AI外呼 + 法律话术 + 合规体系 + 微信跟进」捏成闭环', 14, WHITE),
    ('通用外呼厂商不做行业纵深，法律科技厂商不做案源获取 —— 中间是真空地带', 14, GRAY_TEXT),
    ('但「没人做」也可能是「不值得做」—— 详见反向论证', 14, RED),
])

add_multiline(s, Inches(8.5), Inches(4.7), Inches(4.3), Inches(2.0), [
    ('⚠ 竞品「不做」的另一种解释', 14, ORANGE, True),
    ('', 6, GRAY_TEXT),
    ('1. 有付费意愿 + 出得起钱的律所可能<1000家', 11, GRAY_TEXT),
    ('2. 律所获客成本高（靠地推，人均月拜访20家成交1-2家）', 11, GRAY_TEXT),
    ('3. 法律行业合规雷太大，小公司扛不住', 11, GRAY_TEXT),
    ('4. 案由多样，话术定制需求碎，投入产出比差', 11, GRAY_TEXT),
])

add_source_footer(s, '竞品信息：各厂商官网公开定价页 · 法律科技行业分析报告（2025）· 阿里云/腾讯云语音服务官网')


# ════════════════════════════════════════════
# SLIDE 5 — 差异化价值
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5), '差异化价值：三个可防御壁垒', font_size=28, color=WHITE, bold=True)

barriers = [
    ('话术即产品', CYAN, [
        '3套专业法律话术（离婚/债务/交通事故），每套五环节设计',
        '不是「附赠说明书」，是独立收费产品。可扩展为「按案由话术库」',
        '客户回应应对话术表（6种常见回应 → 标准化应答）',
        '关键数据：接通率>60% / 意向率>20% / 加微信率>50% / 到店率>30%',
    ]),
    ('合规即交付物', GREEN, [
        '10大类 60+ 合规检查项，覆盖个保法/通信外呼/广告法',
        '自动生成合规报告 + 承诺书模板 → 律所最恐惧的法律风险被化解',
        '沟通成本优势：销售时可以拿合规清单直接打消主任的「骚扰投诉」顾虑',
        '关键事件：个人信息保护法违规罚款最高5000万元或营业额5%',
    ]),
    ('闭环数据模型', PURPLE, [
        '接通→意向→微信→到店→委托 五级转化漏斗，每个环节可量化',
        '律所主任首次可以看到「哪条渠道亏了哪条赚了」',
        '转化数据越积累，话术优化越精准 → 网络效应虽弱但有数据飞轮',
        '关键指标：1个律所5人团队，月节省≥¥10,000（替代1个电销人员）',
    ]),
]

for i, (title, accent, items) in enumerate(barriers):
    y = Inches(1.0 + i * 2.1)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.9))
    add_accent_bar(s, Inches(0.5), y, Inches(0.06), Inches(1.9), accent)
    add_textbox(s, Inches(0.9), y + Inches(0.1), Inches(3), Inches(0.4),
                title, font_size=20, color=accent, bold=True)
    add_multiline(s, Inches(0.9), y + Inches(0.6), Inches(11.5), Inches(1.2),
                  [(f'• {item}', 13, GRAY_TEXT) for item in items])

add_source_footer(s, '话术来源：项目文件 docs/SCRIPTS.md（3套完整法律话术）· 合规来源：项目文件 docs/COMPLIANCE.md（10类60+检查项）· 个人信息保护法第66条')


# ════════════════════════════════════════════
# SLIDE 6 — 反向论证（危险假设）
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5), '反向论证：假设它会失败，拆解致命漏洞', font_size=28, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(0.65), Inches(10), Inches(0.35),
            '任何商业计划书都需要回答：什么情况下这个事成不了？', font_size=14, color=RED)

fatal_flaws = [
    ('假设 1', '「律所缺案源就愿意为工具付费」', '可能有案源焦虑+出得起钱+不抵触技术的律所<1000家。律师天性怀疑一切、精算每笔支出，培育SaaS付费心智需3-5年而非1-2年', RED),
    ('假设 2', '「外呼这条赛道还在」', '手机管家陌生号码拦截默认开启。工信部2023年起整治趋严。外呼号码被标记骚扰后接通率<20%。监管和市场在双重绞杀外呼模式', RED),
    ('假设 3', '「AI打电话转化率≥人工」', 'Demo里意向等级是用random.random()随机生成的。实测ASR在电话信道上的识别率可能<70%。TTS机械感让客户秒识破。一旦识破，信任归零', RED),
    ('假设 4', '「200客户能盈利」', '漏算了：语音线路月租¥3-5K、客户成功人员¥12K、OS存储录音¥1-3K、投诉处理成本不可估。实际月成本可能¥60-70K，200客户的¥66K刚刚打平', ORANGE),
    ('假设 5', '「法律行业合规风险可控」', '一个客户投诉「AI冒充律师」，面临的是个保法最高5000万罚款。LLM在法律场景一句话说错就是「非法执业」。这不是技术问题，是刑事责任问题', RED),
]

for i, (tag, assumption, risk, color) in enumerate(fatal_flaws):
    y = Inches(1.15 + i * 1.2)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.08))
    add_accent_bar(s, Inches(0.5), y, Inches(0.06), Inches(1.08), color)
    add_textbox(s, Inches(0.9), y + Inches(0.08), Inches(1.5), Inches(0.3),
                tag, font_size=12, color=color, bold=True)
    add_textbox(s, Inches(2.3), y + Inches(0.08), Inches(4.5), Inches(0.3),
                assumption, font_size=14, color=WHITE, bold=True)
    add_textbox(s, Inches(0.9), y + Inches(0.5), Inches(11.5), Inches(0.5),
                risk, font_size=12, color=GRAY_TEXT)

add_source_footer(s, '监管来源：工信部《通信短信息服务管理规定》2023修订版 · 个人信息保护法 · Apple/华为手机陌生号码拦截机制公开文档')


# ════════════════════════════════════════════
# SLIDE 7 — 技术真难点
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5), '技术真难点：语音机器人如何不被识别出来？', font_size=28, color=WHITE, bold=True)

# Problem statement
add_rect(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.65), fill=RGBColor(0x1A, 0x10, 0x18))
add_textbox(s, Inches(0.8), Inches(0.92), Inches(11.7), Inches(0.5),
            '核心命题：客户什么时候发现对面不是人？一旦发现，信任归零。律师场景比快递/外卖严重——客户在聊离婚、债务、事故，情绪本已脆弱，发现是机器，辱骂举报拉黑一气呵成。',
            font_size=14, color=RED)

exposure_points = [
    ('前3秒\n音色不对', '人3秒内判断对方的年龄、情绪、态度。TTS的韵律、中文轻声、儿化、语流音变在自然度上仍有鸿沟。2025年CosyVoice 2 / ChatTTS已有显著突破，但电话信道8kHz掩盖部分合成痕迹，反而有帮助。', '中等', ORANGE),
    ('第一次卡顿\n延迟异常', '真人说+听+想的循环是不规则的。ASR→LLM→TTS链路的响应延迟与人的200ms自然响应有差距，且停顿时长一致性暴露机器特性。需要加入非对称延迟和随机抖动来模拟思考。', '致命', RED),
    ('客户打断\nbarge-in', '这是最容易暴露的节点——客户说「等一下——」，机器人继续念词不停。barge-in需要流式ASR+流式LLM+流式TTS+打断策略联动，不是某一个模型的事。目前工程界整体方案都不成熟。', '致命', RED),
    ('情绪错位\n共情失败', '客户哭了，机器人还是标准热情腔。TTS情感控制只能「预设一种风格」，无法「跟随上下文动态变化」。而法律咨询场景天然高情绪。共情是先于内容的——人能在两句之内感知到对方理解自己的情绪。', '严重', RED),
    ('「你是机器人吧」\n直接问', '一旦被问，LLM撒谎（「不是，我是真人」）违规，承认则报废。两难。可能的解法是设计一套「幽默化解+立即转人工」的策略，但不能依靠LLM即兴发挥。', '严重', RED),
]

for i, (phase, detail, severity, color) in enumerate(exposure_points):
    y = Inches(1.65 + i * 1.1)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0))
    add_rect(s, Inches(0.5), y, Inches(0.06), Inches(1.0), color)
    # Phase
    add_textbox(s, Inches(0.8), y + Inches(0.08), Inches(1.5), Inches(0.85),
                phase, font_size=11, color=color, bold=True)
    # Detail
    add_textbox(s, Inches(2.5), y + Inches(0.08), Inches(8.5), Inches(0.85),
                detail, font_size=12, color=GRAY_TEXT)
    # Severity tag
    tag_color = RED if severity == '致命' else ORANGE
    add_textbox(s, Inches(11.3), y + Inches(0.3), Inches(1.3), Inches(0.35),
                severity, font_size=11, color=tag_color, bold=True, alignment=PP_ALIGN.CENTER)

add_source_footer(s, 'TTS技术现状：CosyVoice 2 (阿里2024) · ChatTTS (2024开源) · Fish Audio (2025) · 火山引擎语音合成API文档。ASR：阿里云NLS / 讯飞实时语音识别 · Whisper')


# ════════════════════════════════════════════
# SLIDE 8 — 商业模型：乐观 vs 悲观
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5), '商业模型：乐观测算与悲观现实', font_size=28, color=WHITE, bold=True)

# ── Optimistic ──
add_rect(s, Inches(0.5), Inches(0.85), Inches(6.0), Inches(6.0))
add_accent_bar(s, Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.06), GREEN)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.4), '乐观测算', font_size=18, color=GREEN, bold=True)

optimistic_data = [
    ('定价方案', [
        '基础版：¥1,980/月（50通/天，1套话术）',
        '标准版：¥3,980/月（100通/天，5套话术+微信）',
        '企业版：¥6,980/月（300通/天，话术定制+合规报告）',
        '通话费附加：¥0.36/通（阿里云成本¥0.15，毛利¥0.21）',
    ]),
    ('规模化测算（标准版为主）', [
        '50客户：年收入 ¥180万 + 通话毛利 ¥19万 = ¥199万',
        '200客户：年收入 ¥720万 + 通话毛利 ¥76万 = ¥796万',
        '500客户：年收入¥1800万 + 通话毛利¥189万 =¥1,989万',
        '1,000客户：年收入¥3360万 + 通话毛利¥378万 =¥3,738万',
    ]),
    ('估值参考', [
        '种子期（50-100客户）：¥1,000-3,200万（5-8x ARR）',
        'A轮（200-500客户）：¥5,000万-2亿（6-10x ARR）',
    ]),
]

y_off = Inches(1.5)
for section_title, items in optimistic_data:
    add_textbox(s, Inches(0.8), y_off, Inches(5.5), Inches(0.3),
                section_title, font_size=13, color=CYAN, bold=True)
    y_off += Inches(0.3)
    for item in items:
        add_textbox(s, Inches(1.1), y_off, Inches(5.2), Inches(0.25),
                    f'• {item}', font_size=11, color=GRAY_TEXT)
        y_off += Inches(0.25)
    y_off += Inches(0.15)

# ── Pessimistic ──
add_rect(s, Inches(6.8), Inches(0.85), Inches(6.0), Inches(6.0))
add_accent_bar(s, Inches(6.8), Inches(0.85), Inches(6.0), Inches(0.06), RED)
add_textbox(s, Inches(7.1), Inches(1.0), Inches(5.5), Inches(0.4), '悲观现实', font_size=18, color=RED, bold=True)

pessimistic_data = [
    ('隐性成本（月度）', [
        '云服务（ECS+RDS+Redis）：¥3,000',
        '阿里云语音（ASR+TTS+线路）：¥18,000-20,000',
        'LLM API（硅基流动/MiniMax）：¥5,000',
        '客户成功1人：¥12,000',
        '录音OSS存储（合规6个月+）：¥1,000-3,000',
        '研发1人：¥15,000',
        '投诉处理/合规/法务：不可估',
        '月总成本实际：¥60,000-70,000 ⚠',
    ]),
    ('200客户时', [
        '月收入：¥66,000',
        '月成本：¥60,000-70,000',
        '净利润率：~0-10%',
        '「赚钱」前提是所有假设同时成立',
    ]),
    ('漏算的关键成本', [
        '阿里云线路月租 ¥3,000-5,000',
        '律所不会自学，onboarding 成本高',
        '地推获客成本：一家律所 ¥3,000-8,000',
        '合规事故准备金：建议预留年收入10-20%',
    ]),
]

y_off = Inches(1.5)
for section_title, items in pessimistic_data:
    add_textbox(s, Inches(7.1), y_off, Inches(5.5), Inches(0.3),
                section_title, font_size=13, color=ORANGE, bold=True)
    y_off += Inches(0.3)
    for item in items:
        add_textbox(s, Inches(7.4), y_off, Inches(5.2), Inches(0.25),
                    f'• {item}', font_size=11, color=GRAY_TEXT)
        y_off += Inches(0.25)
    y_off += Inches(0.15)

add_source_footer(s, '价格参考：百应科技/探迹/励销云公开定价 · 阿里云语音服务语音通话定价页（0.15元/通起）· 云服务器ECS 2核4G月费 · 法律科技SaaS行业薪酬报告')


# ════════════════════════════════════════════
# SLIDE 9 — 市场 vs 监管
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5), '市场 vs 监管：外呼赛道的生死线', font_size=28, color=WHITE, bold=True)

# Left — market opportunity
add_rect(s, Inches(0.5), Inches(0.9), Inches(6.0), Inches(2.8))
add_accent_bar(s, Inches(0.5), Inches(0.9), Inches(6.0), Inches(0.06), CYAN)
add_textbox(s, Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.3), '市场端 — 可触达客户', font_size=16, color=CYAN, bold=True)
add_multiline(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(2.2), [
    ('TAM (总可触达市场): ~4万家中国律所', 14, WHITE, True),
    ('', 4, GRAY_TEXT),
    ('SAM (可服务市场): 有案源焦虑 + 10人以下 + 年营收>50万', 13, GRAY_TEXT),
    ('  ≈ 5,000-8,000 家（乐观估计头部15-25%）', 13, GRAY_TEXT),
    ('', 4, GRAY_TEXT),
    ('SOM (可获取市场): 舍得为工具付费 + 不抵触技术', 13, GRAY_TEXT),
    ('  ≈ 500-2,000 家（悲观估计头部1-5%）', 13, RED, True),
    ('', 4, GRAY_TEXT),
    ('付费心智：律所过去20年没有「软件订阅」这个品类', 12, ORANGE),
    ('培育时间可能需要 3-5 年，而非 1-2 年', 12, ORANGE),
])

# Right — regulatory headwind
add_rect(s, Inches(6.8), Inches(0.9), Inches(6.0), Inches(2.8))
add_accent_bar(s, Inches(6.8), Inches(0.9), Inches(6.0), Inches(0.06), RED)
add_textbox(s, Inches(7.1), Inches(1.0), Inches(5.5), Inches(0.3), '监管端 — 绞杀力量', font_size=16, color=RED, bold=True)
add_multiline(s, Inches(7.1), Inches(1.4), Inches(5.5), Inches(2.2), [
    ('工信部 2023：商业营销电话整治趋严', 14, RED, True),
    ('未经同意的商业外呼明确违规', 13, GRAY_TEXT),
    ('', 4, GRAY_TEXT),
    ('手机生态：Apple/华为/小米默认拦截陌生号码', 14, RED, True),
    ('被标记为骚扰电话后接通率 <20%', 13, GRAY_TEXT),
    ('', 4, GRAY_TEXT),
    ('个人信息保护法第66条', 14, RED, True),
    ('罚款：最高 5,000 万元或上年营业额 5%', 13, GRAY_TEXT),
    ('', 4, GRAY_TEXT),
    ('95/96 号码门槛：企业资质+月租¥3,000+', 13, ORANGE),
    ('中小律所基本申请不下来', 13, ORANGE),
])

# Bottom — The blunt truth
add_rect(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(2.1), fill=RGBColor(0x1A, 0x08, 0x08))
add_textbox(s, Inches(0.8), Inches(4.15), Inches(11.7), Inches(0.4),
            '核心矛盾', font_size=18, color=RED, bold=True)
add_multiline(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.4), [
    ('一个产品不能建立在一个正在被监管和市场双重绞杀的通信方式上。', 16, WHITE, True),
    ('', 6, GRAY_TEXT),
    ('就算话术再好、AI 再聪明：', 14, GRAY_TEXT),
    ('— 如果电话打不通（手机自动拦截），等于零', 14, GRAY_TEXT),
    ('— 如果打通的都被标记骚扰，等于负资产', 14, GRAY_TEXT),
    ('— 如果因为一个投诉被罚 5,000 万，公司直接死亡', 14, GRAY_TEXT),
    ('', 6, GRAY_TEXT),
    ('可能的出路：从「主动外呼」转向「先微信建联 → 客户同意后再电话」，改变触点发生顺序。但这也意味着系统核心要从外呼引擎变为微信对话引擎。', 14, CYAN),
])

add_source_footer(s, '监管来源：工信部《通信短信息服务管理规定》（2023修订）· 个人信息保护法 · Apple iOS 陌生号码静音功能文档 · 华为EMUI骚扰拦截机制 · 阿里云语音服务号码申请条件')


# ════════════════════════════════════════════
# SLIDE 10 — 验证路径
# ════════════════════════════════════════════
s = new_slide()
add_accent_bar(s, Inches(0), Inches(0), Inches(0.08), Inches(0.6))
add_textbox(s, Inches(0.8), Inches(0.15), Inches(10), Inches(0.5), '下一步：验证最危险的假设', font_size=28, color=WHITE, bold=True)

stages = [
    ('阶段一：核心验证 (1-2周)', CYAN, [
        'P0 | 轮换所有已暴露密钥（钥匙.md），确保安全基线',
        'P0 | 配置阿里云语音服务，用已写好的aliyun-call.service.ts打通一通实际电话',
        'P1 | 用10个真实号码（非随机生成）跑一轮完整外呼→ASR→LLM→TTS',
        '判断标准：AI能主导完成一通2分钟以上的对话，客户未察觉是机器人',
    ]),
    ('阶段二：真实测试 (2-3周)', GREEN, [
        '对接1家律所种子客户（你认识的事务所），用真实客户名单跑一轮',
        '收集四率数据：接通率 / 意向率 / 加微信率 / 到店率',
        '拿数据去问5个律所主任——「这个效果，你愿意每月付多少钱？」',
        '判断标准：≥3个主任给出明确付费金额，且金额覆盖成本',
    ]),
    ('阶段三：最小产品 (3-4周)', PURPLE, [
        '前端只做三个页面：客户导入 → 外呼任务 → 通话结果列表',
        '加入黑名单+频率限制（之前已完成90%）',
        '不做多租户、不做高级报表、不做微信集成——先卖出去再说',
        '判断标准：第一个付费客户上线，哪怕只收¥500/月',
    ]),
]

for i, (title, accent, items) in enumerate(stages):
    y = Inches(0.85 + i * 2.15)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(2.0))
    add_accent_bar(s, Inches(0.5), y, Inches(0.06), Inches(2.0), accent)
    add_textbox(s, Inches(0.9), y + Inches(0.1), Inches(3.5), Inches(0.35),
                title, font_size=18, color=accent, bold=True)
    add_multiline(s, Inches(0.9), y + Inches(0.55), Inches(11.5), Inches(1.3),
                  [(f'{item}', 13, GRAY_TEXT) for item in items])

add_source_footer(s, '验证方法论：Steve Blank《The Four Steps to the Epiphany》客户开发模型 · Y Combinator 「Do Things That Don\'t Scale」原则')


# ════════════════════════════════════════════
# SLIDE 11 — 总结
# ════════════════════════════════════════════
s = new_slide()
add_rect(s, Inches(0), Inches(0), W, H, fill=BG_DARK)
add_accent_bar(s, Inches(0), Inches(3.0), Inches(0.1), Inches(1.2), CYAN)
add_textbox(s, Inches(1), Inches(1.8), Inches(11), Inches(1.0),
            '总结', font_size=48, color=WHITE, bold=True)

add_multiline(s, Inches(1), Inches(3.5), Inches(11), Inches(3.0), [
    ('方向判断：选对了痛点（案源是律所第一焦虑），打对了差异化（通用外呼不做法务，法务科技不做外呼），但外呼这条赛道本身在塌方。', 16, GRAY_TEXT),
    ('', 10, GRAY_TEXT),
    ('产品判断：技术框架完成度高（85%），但「完成度」和「可用性」之间隔着一通实际电话。Demo里意向等级是random.random()随机出来的，不是ASR从真实客户语音转出来的。', 16, GRAY_TEXT),
    ('', 10, GRAY_TEXT),
    ('商业判断：算得过账——律所每月花¥2,000替代一个¥6,000底薪的电销，ROI不需要解释。但算得过账≠有人付钱，律所的SaaS付费心智可能是整个故事里最被低估的阻力。', 16, GRAY_TEXT),
    ('', 10, GRAY_TEXT),
    ('技术判断：TTS不是瓶颈，对话体验的不可识别性才是。暴露节点已拆解（延迟、打断、情绪、被质问），每个都比「声音不够像真人」致命得多。', 16, GRAY_TEXT),
    ('', 10, GRAY_TEXT),
    ('最后一步：不是写更多代码，是用阿里云语音服务打一通实际电话，录下来听效果。这比任何商业计划书都有说服力。', 16, CYAN, True),
])

# Bottom bar
add_accent_bar(s, Inches(0), H - Inches(0.03), W, Inches(0.03), CYAN)

# Save
output_path = '/Users/apple/Documents/智能外呼系统/智能外呼系统_商业分析.pptx'
prs.save(output_path)
print(f'PPT saved to: {output_path}')
print(f'Slides: {len(prs.slides)}')
